import streamlit as st
import pandas as pd
from pptx import Presentation
import io
import json
import urllib.request
from openai import OpenAI
from openpyxl import Workbook
from openpyxl.styles import PatternFill, Font, Alignment, Border, Side
from openpyxl.utils import get_column_letter
import datetime

st.set_page_config(
    page_title="Déroulé Pédagogique Qualiopi",
    page_icon="📋",
    layout="wide"
)

TEAL     = "#4EC5C0"
TEAL_MID = "#A8E0DD"
TEAL_LIGHT = "#EBF7F7"

COLS = [
    "Jour",
    "Horaires",
    "Objectifs pédagogiques",
    "Contenu de la séquence",
    "Moyens pédagogiques",
    "Modalités de validation des acquis",
]

LOGO_URL = "https://socentric.fr/wp-content/uploads/2022/09/cropped-SoCentric_Logo-01.png"


@st.cache_data(show_spinner=False)
def load_logo():
    try:
        return urllib.request.urlopen(LOGO_URL, timeout=10).read()
    except Exception:
        return None


# ─────────────────────────────────────────────
# Helpers
# ─────────────────────────────────────────────

def get_api_key():
    try:
        return st.secrets.get("OPENAI_API_KEY", "")
    except Exception:
        return ""


def check_password():
    if st.session_state.get("authenticated"):
        return True
    st.markdown("## 🔒 Accès protégé")
    pwd = st.text_input("Mot de passe", type="password", placeholder="Entrez le mot de passe…")
    if st.button("Se connecter", type="primary"):
        try:
            correct = st.secrets.get("APP_PASSWORD", "")
        except Exception:
            correct = ""
        if pwd and pwd == correct:
            st.session_state["authenticated"] = True
            st.rerun()
        else:
            st.error("Mot de passe incorrect.")
    st.stop()


def extract_pptx_text(pptx_bytes: bytes) -> str:
    prs = Presentation(io.BytesIO(pptx_bytes))
    slides = []
    for i, slide in enumerate(prs.slides):
        texts = [
            shape.text.strip()
            for shape in slide.shapes
            if hasattr(shape, "text") and shape.text.strip()
        ]
        if texts:
            slides.append(f"=== Slide {i + 1} ===\n" + "\n".join(texts))
    return "\n\n".join(slides)


def analyze_with_gpt(pptx_text: str, sequences_fixes: str, nb_jours: int, api_key: str) -> list[dict]:
    client = OpenAI(api_key=api_key)
    prompt = f"""Tu es un expert en ingénierie de formation et certifications Qualiopi.

Tu dois construire un déroulé pédagogique complet sur {nb_jours} jour(s) en combinant deux sources :

---
SOURCE 1 — SÉQUENCES FIXES (toujours présentes, quelle que soit la formation) :
{sequences_fixes}
---
SOURCE 2 — CONTENU PÉDAGOGIQUE (extrait du support de formation) :
{pptx_text}
---

Ta mission :
1. Intègre TOUTES les séquences fixes aux bons moments (accueil en début, pauses matin/après-midi, déjeuner à 12h30, ice breaker en début d'après-midi, positionnements en début/fin de journée, satisfaction en fin de journée, etc.).
2. Regroupe le contenu pédagogique en séquences logiques et insère-les entre les séquences fixes.
3. Répartis les horaires sur {nb_jours} jour(s), journée 9h00 → 17h00 (1h déjeuner à 12h30, 2 pauses de 15 min).
4. Pour les séquences fixes : objectifs/moyens/modalités courts et adaptés.
5. Pour les séquences contenu : objectifs avec verbe d'action, moyens variés, modalités précises.

Réponds UNIQUEMENT avec un tableau JSON valide (sans markdown, sans texte autour) :
[
  {{
    "jour": "J1",
    "horaires": "09h00 - 09h30",
    "objectifs_pedagogiques": "...",
    "contenu_sequence": "...",
    "moyens_pedagogiques": "...",
    "modalites_validation": "..."
  }}
]"""

    response = client.chat.completions.create(
        model="gpt-4o-mini",
        messages=[{"role": "user", "content": prompt}],
        temperature=0.3,
        max_tokens=4000,
    )
    raw = response.choices[0].message.content.strip()
    if raw.startswith("```"):
        raw = raw.split("```")[1]
        if raw.startswith("json"):
            raw = raw[4:]
        raw = raw.strip()

    return [
        {
            "Jour":                               item.get("jour", ""),
            "Horaires":                           item.get("horaires", ""),
            "Objectifs pédagogiques":             item.get("objectifs_pedagogiques", ""),
            "Contenu de la séquence":             item.get("contenu_sequence", ""),
            "Moyens pédagogiques":                item.get("moyens_pedagogiques", ""),
            "Modalités de validation des acquis": item.get("modalites_validation", ""),
        }
        for item in json.loads(raw)
    ]


# ─────────────────────────────────────────────
# Export XLSX
# ─────────────────────────────────────────────

def generate_excel(meta: dict, df: pd.DataFrame) -> bytes:
    wb = Workbook()
    ws = wb.active
    ws.title = "Déroulé Pédagogique"

    f_teal  = PatternFill("solid", fgColor=TEAL.lstrip("#"))
    f_light = PatternFill("solid", fgColor=TEAL_LIGHT.lstrip("#"))
    f_white = PatternFill("solid", fgColor="FFFFFF")

    thin   = Side(border_style="thin",   color="AAAAAA")
    medium = Side(border_style="medium", color=TEAL.lstrip("#"))
    bdr    = Border(left=thin,   right=thin,   top=thin,   bottom=thin)
    bdr_h  = Border(left=medium, right=medium, top=medium, bottom=medium)

    center = Alignment(horizontal="center", vertical="center", wrap_text=True)
    left   = Alignment(horizontal="left",   vertical="top",    wrap_text=True)

    # Titre
    ws.merge_cells("A1:F1")
    c = ws["A1"]
    c.value     = f"DÉROULÉ PÉDAGOGIQUE — {meta['nom_formation'].upper()}"
    c.font      = Font(name="Calibri", bold=True, color="FFFFFF", size=14)
    c.fill      = f_teal
    c.alignment = center
    ws.row_dimensions[1].height = 32

    # Métadonnées
    for rng, val in [
        ("A2:B2", f"Formateur : {meta['formateur']}"),
        ("C2:D2", f"Dates : {meta['dates']}"),
        ("E2:F2", f"Lieu : {meta['lieu']}  |  Stagiaires : {meta['nb_stagiaires']}"),
    ]:
        ws.merge_cells(rng)
        c = ws[rng.split(":")[0]]
        c.value     = val
        c.font      = Font(name="Calibri", color="FFFFFF", size=10)
        c.fill      = PatternFill("solid", fgColor=TEAL_MID.lstrip("#"))
        c.alignment = center
    ws.row_dimensions[2].height = 22

    # Rappel des objectifs
    next_row = 3
    if meta.get("rappel_objectifs"):
        ws.merge_cells("A3:F3")
        c = ws["A3"]
        c.value     = f"Rappel des objectifs : {meta['rappel_objectifs']}"
        c.font      = Font(name="Calibri", bold=True, italic=True, size=10, color=TEAL.lstrip("#"))
        c.alignment = left
        c.border    = bdr
        ws.row_dimensions[3].height = 35
        next_row = 4

    # En-têtes
    headers = [
        ("JOUR", 8),
        ("HORAIRES", 13),
        ("OBJECTIFS PÉDAGOGIQUES", 32),
        ("CONTENU DE LA SÉQUENCE", 44),
        ("MOYENS PÉDAGOGIQUES\nComment se fait la transmission ?", 30),
        ("MODALITÉS DE VALIDATION DES ACQUIS\nComment évaluer l'atteinte de l'objectif ?", 32),
    ]
    hr = next_row
    for ci, (label, width) in enumerate(headers, 1):
        c = ws.cell(row=hr, column=ci, value=label)
        c.font      = Font(name="Calibri", bold=True, color="FFFFFF", size=10)
        c.fill      = f_teal
        c.alignment = center
        c.border    = bdr_h
        ws.column_dimensions[get_column_letter(ci)].width = width
    ws.row_dimensions[hr].height = 45

    # Données
    for ri, (_, row) in enumerate(df.iterrows()):
        er   = hr + 1 + ri
        fill = f_light if ri % 2 == 0 else f_white
        for ci, col in enumerate(COLS, 1):
            c = ws.cell(row=er, column=ci, value=str(row.get(col, "") or ""))
            c.font      = Font(name="Calibri", size=10)
            c.fill      = fill
            c.alignment = left
            c.border    = bdr
        ws.row_dimensions[er].height = 65

    # Pied de page (dans une ligne sous le tableau)
    footer_row = hr + 1 + len(df)
    ws.merge_cells(f"A{footer_row}:D{footer_row}")
    c = ws[f"A{footer_row}"]
    c.value = f"{meta.get('version','V1')} ({meta['nom_formation']}) • mis à jour le {meta['date_maj']}"
    c.font  = Font(name="Calibri", italic=True, size=8, color="888888")
    c.alignment = left
    ws.merge_cells(f"E{footer_row}:F{footer_row}")
    c = ws[f"E{footer_row}"]
    c.value     = "Page 1 / 1"
    c.font      = Font(name="Calibri", size=8, color="888888")
    c.alignment = Alignment(horizontal="right")
    ws.row_dimensions[footer_row].height = 18

    out = io.BytesIO()
    wb.save(out)
    out.seek(0)
    return out.getvalue()


# ─────────────────────────────────────────────
# Export PDF
# ─────────────────────────────────────────────

def generate_pdf(meta: dict, df: pd.DataFrame) -> bytes:
    from reportlab.lib.pagesizes import A4
    from reportlab.lib import colors
    from reportlab.lib.units import cm
    from reportlab.platypus import (
        SimpleDocTemplate, Table, TableStyle,
        Paragraph, Spacer, Image, HRFlowable,
    )
    from reportlab.lib.styles import ParagraphStyle
    from reportlab.lib.enums import TA_RIGHT, TA_CENTER, TA_LEFT
    from reportlab.pdfgen import canvas as rl_canvas

    W, H = A4
    TEAL_C   = colors.HexColor(TEAL)
    WHITE_C  = colors.white
    DARK_C   = colors.HexColor("#2C2C2C")
    GREY_C   = colors.HexColor("#888888")
    LIGHT_C  = colors.HexColor(TEAL_LIGHT)

    def ps(name, **kw):
        return ParagraphStyle(name, **kw)

    s_doc_title = ps("DocTitle", fontSize=11, fontName="Helvetica-Bold",
                     alignment=TA_RIGHT, textColor=DARK_C)
    s_formation = ps("Formation", fontSize=13, fontName="Helvetica-Bold",
                     alignment=TA_RIGHT, textColor=DARK_C)
    s_obj       = ps("Obj", fontSize=9, fontName="Helvetica-Bold", textColor=TEAL_C)
    s_hdr       = ps("Hdr", fontSize=7.5, fontName="Helvetica-Bold",
                     textColor=WHITE_C, alignment=TA_CENTER, leading=10)
    s_cell      = ps("Cell", fontSize=8, fontName="Helvetica",
                     textColor=DARK_C, leading=10)
    s_cell_c    = ps("CellC", fontSize=8, fontName="Helvetica-Bold",
                     textColor=DARK_C, leading=10, alignment=TA_CENTER)

    # ── NumberedCanvas for Page X / Y ──────────────────────
    footer_left = (
        f"{meta.get('version','V1')} ({meta['nom_formation']}) "
        f"• mis à jour le {meta['date_maj']}"
    )

    class NumberedCanvas(rl_canvas.Canvas):
        def __init__(self, *args, **kwargs):
            rl_canvas.Canvas.__init__(self, *args, **kwargs)
            self._saved = []

        def showPage(self):
            self._saved.append(dict(self.__dict__))
            self._startPage()

        def save(self):
            n = len(self._saved)
            for state in self._saved:
                self.__dict__.update(state)
                self.setFont("Helvetica-Oblique", 7)
                self.setFillColor(GREY_C)
                self.drawString(1.5 * cm, 0.75 * cm, footer_left)
                self.drawRightString(W - 1.5 * cm, 0.75 * cm,
                                     f"Page {self._pageNumber} / {n}")
                rl_canvas.Canvas.showPage(self)
            rl_canvas.Canvas.save(self)

    # ── Document ────────────────────────────────────────────
    buf = io.BytesIO()
    doc = SimpleDocTemplate(
        buf, pagesize=A4,
        leftMargin=1.5*cm, rightMargin=1.5*cm,
        topMargin=1.5*cm, bottomMargin=2*cm,
    )

    story = []

    # ── Logo + titre document ────────────────────────────────
    logo_bytes = load_logo()
    if logo_bytes:
        logo = Image(io.BytesIO(logo_bytes), width=3.5*cm, height=1.8*cm)
    else:
        logo = Spacer(3.5*cm, 1.8*cm)

    hdr_tbl = Table(
        [[logo, Paragraph("Déroulé pédagogique", s_doc_title)]],
        colWidths=[9*cm, 8.2*cm],
    )
    hdr_tbl.setStyle(TableStyle([
        ("VALIGN",       (0, 0), (-1, -1), "TOP"),
        ("LEFTPADDING",  (0, 0), (-1, -1), 0),
        ("RIGHTPADDING", (0, 0), (-1, -1), 0),
        ("TOPPADDING",   (0, 0), (-1, -1), 0),
        ("BOTTOMPADDING",(0, 0), (-1, -1), 0),
    ]))
    story.append(hdr_tbl)
    story.append(Spacer(1, 0.5*cm))

    # ── Nom de formation ────────────────────────────────────
    story.append(Paragraph(f"« {meta['nom_formation']} »", s_formation))
    story.append(Spacer(1, 0.4*cm))

    # ── Rappel des objectifs ─────────────────────────────────
    if meta.get("rappel_objectifs"):
        obj_tbl = Table(
            [["", Paragraph(
                f"<b>Rappel des objectifs :</b>  {meta['rappel_objectifs']}",
                s_obj
            )]],
            colWidths=[0.25*cm, 16.95*cm],
        )
        obj_tbl.setStyle(TableStyle([
            ("BACKGROUND",    (0, 0), (0, 0), TEAL_C),
            ("LEFTPADDING",   (1, 0), (1, 0), 6),
            ("TOPPADDING",    (0, 0), (-1, -1), 3),
            ("BOTTOMPADDING", (0, 0), (-1, -1), 3),
        ]))
        story.append(obj_tbl)
        story.append(HRFlowable(width="100%", thickness=1, color=TEAL_C, spaceAfter=4))

    story.append(Spacer(1, 0.25*cm))

    # ── Tableau principal ────────────────────────────────────
    col_w = [1.2*cm, 2.1*cm, 3.4*cm, 5.6*cm, 2.75*cm, 2.75*cm]  # = 17.8cm

    def hdr_p(line1, line2=None):
        txt = f"<b>{line1}</b>"
        if line2:
            txt += f"<br/><font size='6'>{line2}</font>"
        return Paragraph(txt, s_hdr)

    thead = [
        hdr_p("JOUR"),
        hdr_p("HORAIRES"),
        hdr_p("OBJECTIFS", "PÉDAGOGIQUES"),
        hdr_p("CONTENU DE LA SÉQUENCE"),
        hdr_p("MOYENS PÉDAGOGIQUES MIS EN ŒUVRE",
              "Comment se fait la transmission ?"),
        hdr_p("MODALITES DE VALIDATION DES ACQUIS",
              "Comment évaluer l'atteinte de l'objectif ?"),
    ]

    rows = [thead]
    for _, row in df.iterrows():
        rows.append([
            Paragraph(str(row.get("Jour",     "") or ""), s_cell_c),
            Paragraph(str(row.get("Horaires", "") or ""), s_cell_c),
            Paragraph(str(row.get("Objectifs pédagogiques",             "") or ""), s_cell),
            Paragraph(str(row.get("Contenu de la séquence",             "") or ""), s_cell),
            Paragraph(str(row.get("Moyens pédagogiques",                "") or ""), s_cell),
            Paragraph(str(row.get("Modalités de validation des acquis", "") or ""), s_cell),
        ])

    style_cmds = [
        ("BACKGROUND",    (0, 0), (-1,  0), TEAL_C),
        ("TEXTCOLOR",     (0, 0), (-1,  0), WHITE_C),
        ("VALIGN",        (0, 0), (-1, -1), "MIDDLE"),
        ("ALIGN",         (0, 1), ( 1, -1), "CENTER"),
        ("GRID",          (0, 0), (-1, -1), 0.5, colors.HexColor("#CCCCCC")),
        ("BOX",           (0, 0), (-1, -1), 1,   TEAL_C),
        ("TOPPADDING",    (0, 0), (-1, -1), 4),
        ("BOTTOMPADDING", (0, 0), (-1, -1), 4),
        ("LEFTPADDING",   (0, 0), (-1, -1), 4),
        ("RIGHTPADDING",  (0, 0), (-1, -1), 4),
        ("MINROWHEIGHT",  (0, 0), (-1,  0), 1.4*cm),
    ]
    # Alternating row backgrounds
    for i in range(1, len(rows)):
        bg = LIGHT_C if i % 2 == 0 else WHITE_C
        style_cmds.append(("BACKGROUND", (0, i), (-1, i), bg))
        style_cmds.append(("MINROWHEIGHT", (0, i), (-1, i), 1.1*cm))

    main_tbl = Table(rows, colWidths=col_w, repeatRows=1)
    main_tbl.setStyle(TableStyle(style_cmds))
    story.append(main_tbl)

    doc.build(story, canvasmaker=NumberedCanvas)
    buf.seek(0)
    return buf.getvalue()


# ─────────────────────────────────────────────
# UI
# ─────────────────────────────────────────────

check_password()

st.title("📋 Déroulé Pédagogique Qualiopi")
st.caption("Importe ton PPTX → l'IA remplit le tableau Qualiopi → exporte en PDF ou Excel.")

# ── Métadonnées ───────────────────────────────
with st.expander("ℹ️ Informations de la formation", expanded=True):
    c1, c2, c3 = st.columns(3)
    with c1:
        nom_formation = st.text_input("Nom de la formation",
                                      placeholder="Ex : Management - Édouard")
        formateur = st.text_input("Formateur", placeholder="Ex : Jean Dupont")
    with c2:
        date_debut = st.date_input("Date de début", value=datetime.date.today())
        date_fin   = st.date_input("Date de fin",   value=datetime.date.today())
        nb_jours   = max(1, (date_fin - date_debut).days + 1)
        st.caption(f"Durée calculée : **{nb_jours} jour(s)**")
    with c3:
        lieu          = st.text_input("Lieu", placeholder="Ex : Paris / Distanciel")
        nb_stagiaires = st.number_input("Nombre de stagiaires", min_value=1, value=10)
        version       = st.text_input("Version du document", value="V1",
                                      help="Ex : V1, V2… apparaît dans le pied de page")
        date_maj      = st.date_input("Mis à jour le", value=datetime.date.today())

    rappel_objectifs = st.text_area(
        "Rappel des objectifs",
        placeholder="Ex : Développer les compétences managériales...",
        height=70,
    )
    sequences_fixes = st.text_area(
        "Séquences fixes (une par ligne)",
        placeholder=(
            "Tour de table de présentation (ouverture)\n"
            "Positionnement sur le sujet (début de journée)\n"
            "Pause de 15 min (matin)\n"
            "Déjeuner d'1H\n"
            "Rupture pédagogique (ice breaker) en début d'am\n"
            "Pause de 15 min (après-midi)\n"
            "Positionnement sur le sujet (fin de journée)\n"
            "Tour de table de satisfaction (fin de journée)"
        ),
        height=150,
        help="Ces séquences seront placées aux bons moments par l'IA.",
    )

st.divider()

# ── Étape 1 : PPTX ───────────────────────────
st.subheader("1 — Importer le support de formation")
uploaded = st.file_uploader("Glisse ton fichier .pptx ici", type=["pptx"],
                             label_visibility="collapsed")

if "pptx_text" not in st.session_state:
    st.session_state["pptx_text"] = ""
if "df" not in st.session_state:
    st.session_state["df"] = pd.DataFrame(columns=COLS)

if uploaded:
    pptx_text = extract_pptx_text(uploaded.read())
    st.session_state["pptx_text"] = pptx_text
    st.success(f"✅ {pptx_text.count('=== Slide ')} slides chargées")

st.divider()

# ── Étape 2 : IA ──────────────────────────────
st.subheader("2 — Générer le déroulé avec l'IA")
api_key     = get_api_key()
can_generate = bool(api_key) and bool(st.session_state["pptx_text"])

if not api_key:
    st.warning("⚠️ Entre ta clé API OpenAI dans la barre latérale.")
elif not st.session_state["pptx_text"]:
    st.info("⬆️ Importe d'abord un fichier PPTX.")

if st.button("✨ Générer le déroulé", type="primary", disabled=not can_generate):
    with st.spinner("Analyse en cours… (10-20 s)"):
        try:
            rows = analyze_with_gpt(
                st.session_state["pptx_text"],
                sequences_fixes,
                nb_jours,
                api_key,
            )
            st.session_state["df"] = pd.DataFrame(rows)
            st.success(f"✅ {len(rows)} séquences générées !")
        except json.JSONDecodeError as e:
            st.error(f"Erreur JSON : {e}")
        except Exception as e:
            st.error(f"Erreur : {e}")

st.divider()

# ── Étape 3 : Tableau ─────────────────────────
st.subheader("3 — Vérifier et ajuster")
edited = st.data_editor(
    st.session_state["df"],
    num_rows="dynamic",
    use_container_width=True,
    height=500,
    column_config={
        "Jour":                               st.column_config.TextColumn(width=70),
        "Horaires":                           st.column_config.TextColumn(width=120),
        "Objectifs pédagogiques":             st.column_config.TextColumn(width=230),
        "Contenu de la séquence":             st.column_config.TextColumn(width=340),
        "Moyens pédagogiques":                st.column_config.TextColumn(width=210),
        "Modalités de validation des acquis": st.column_config.TextColumn(width=210),
    },
)
st.session_state["df"] = edited

st.divider()

# ── Étape 4 : Export ──────────────────────────
st.subheader("4 — Exporter")

def build_meta():
    return {
        "nom_formation":    nom_formation,
        "formateur":        formateur,
        "dates":            f"{date_debut.strftime('%d/%m/%Y')} → {date_fin.strftime('%d/%m/%Y')}",
        "lieu":             lieu,
        "nb_stagiaires":    nb_stagiaires,
        "rappel_objectifs": rappel_objectifs,
        "version":          version,
        "date_maj":         date_maj.strftime("%d/%m/%Y"),
    }

col_pdf, col_xlsx = st.columns(2)

with col_pdf:
    if st.button("📄 Générer le PDF", type="primary", use_container_width=True):
        if not nom_formation.strip():
            st.warning("⚠️ Renseigne le nom de la formation.")
        elif edited.empty:
            st.warning("⚠️ Le tableau est vide.")
        else:
            with st.spinner("Génération du PDF…"):
                try:
                    pdf_bytes = generate_pdf(build_meta(), edited)
                    filename  = (
                        f"Deroulé_Pédagogique_{nom_formation.replace(' ','_')}"
                        f"_{version}_{date_debut.strftime('%Y%m%d')}.pdf"
                    )
                    st.download_button(
                        "💾 Télécharger le PDF",
                        data=pdf_bytes,
                        file_name=filename,
                        mime="application/pdf",
                        type="primary",
                    )
                except Exception as e:
                    st.error(f"Erreur PDF : {e}")

with col_xlsx:
    if st.button("📊 Générer l'Excel", type="secondary", use_container_width=True):
        if not nom_formation.strip():
            st.warning("⚠️ Renseigne le nom de la formation.")
        elif edited.empty:
            st.warning("⚠️ Le tableau est vide.")
        else:
            excel_bytes = generate_excel(build_meta(), edited)
            filename    = (
                f"Deroulé_Pédagogique_{nom_formation.replace(' ','_')}"
                f"_{version}_{date_debut.strftime('%Y%m%d')}.xlsx"
            )
            st.download_button(
                "💾 Télécharger l'Excel",
                data=excel_bytes,
                file_name=filename,
                mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            )
